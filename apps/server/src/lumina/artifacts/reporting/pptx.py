from __future__ import annotations

import textwrap
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.util import Inches, Pt

from .model import ReportDocument


_COBALT = RGBColor(49, 95, 189)
_INK = RGBColor(32, 38, 49)
_MUTED = RGBColor(107, 114, 128)
_LIGHT_BLUE = RGBColor(234, 240, 251)


def generate_pptx(report: ReportDocument) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    _add_title_slide(presentation, report)
    _add_paginated_slides(
        presentation,
        "요약 및 원 요청",
        [report.executive_summary, f"원 요청\n{report.request}"],
    )
    _add_metrics_slide(presentation, report)
    for section in report.sections:
        paragraphs = [section.body] if section.body else []
        paragraphs.extend(f"• {item}" for item in section.bullets)
        _add_paginated_slides(presentation, section.heading, paragraphs)
    if report.action_items:
        _add_paginated_slides(
            presentation,
            "후속 조치",
            [f"{index}. {item}" for index, item in enumerate(report.action_items, 1)],
        )
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _add_title_slide(presentation: PresentationType, report: ReportDocument) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _COBALT
    accent.line.fill.background()
    kicker = slide.shapes.add_textbox(
        Inches(0.85), Inches(1.2), Inches(11.3), Inches(0.35)
    )
    _set_text(kicker.text_frame, "LUMINA · AGENT REPORT", 13, _COBALT, bold=True)
    title = slide.shapes.add_textbox(
        Inches(0.85), Inches(1.75), Inches(11.4), Inches(1.5)
    )
    _set_text(
        title.text_frame, report.title, _title_size(report.title, 34), _INK, bold=True
    )
    summary = slide.shapes.add_textbox(
        Inches(0.9), Inches(3.55), Inches(10.9), Inches(1.35)
    )
    _set_text(summary.text_frame, report.executive_summary, 18, _INK)
    footer = slide.shapes.add_textbox(
        Inches(0.9), Inches(6.65), Inches(11), Inches(0.3)
    )
    _set_text(footer.text_frame, "편집 가능한 Lumina Artifact", 10, _MUTED)


def _add_metrics_slide(presentation: PresentationType, report: ReportDocument) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_slide_title(slide, "핵심 지표")
    count = len(report.metrics)
    gap = 0.25
    width = min(3.75, (11.6 - gap * max(count - 1, 0)) / max(count, 1))
    start = (13.333 - (width * count + gap * max(count - 1, 0))) / 2
    for index, metric in enumerate(report.metrics):
        left = start + index * (width + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(2.0),
            Inches(width),
            Inches(2.65),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _LIGHT_BLUE
        shape.line.color.rgb = _LIGHT_BLUE
        shape.text_frame.clear()
        shape.text_frame.margin_left = Inches(0.25)
        shape.text_frame.margin_right = Inches(0.25)
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        value = shape.text_frame.paragraphs[0]
        value.alignment = PP_ALIGN.CENTER
        run = value.add_run()
        run.text = metric.value
        run.font.name = "맑은 고딕"
        run.font.size = Pt(27)
        run.font.bold = True
        run.font.color.rgb = _COBALT
        label = shape.text_frame.add_paragraph()
        label.alignment = PP_ALIGN.CENTER
        label.space_before = Pt(9)
        run = label.add_run()
        run.text = metric.label
        run.font.name = "맑은 고딕"
        run.font.size = Pt(14)
        run.font.color.rgb = _INK


def _add_body_slide(
    presentation: PresentationType, title: str, paragraphs: list[str]
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_slide_title(slide, title)
    body = slide.shapes.add_textbox(
        Inches(0.85), Inches(1.35), Inches(11.65), Inches(5.35)
    )
    frame = body.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.1)
    frame.margin_right = Inches(0.1)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    for index, text in enumerate(paragraphs or [""]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.font.name = "맑은 고딕"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = _INK
        paragraph.space_after = Pt(10)


def _add_paginated_slides(
    presentation: PresentationType, title: str, paragraphs: list[str]
) -> None:
    pages = _paginate(paragraphs)
    for page_number, page in enumerate(pages, start=1):
        suffix = f" ({page_number}/{len(pages)})" if len(pages) > 1 else ""
        _add_body_slide(presentation, title + suffix, page)


def _add_slide_title(slide: Slide, title: str) -> None:
    title_box = slide.shapes.add_textbox(
        Inches(0.85), Inches(0.42), Inches(11.6), Inches(0.65)
    )
    _set_text(title_box.text_frame, title, _title_size(title, 28), _INK, bold=True)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.85),
        Inches(1.12),
        Inches(1.0),
        Inches(0.06),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _COBALT
    accent.line.fill.background()


def _set_text(
    frame: TextFrame,
    text: str,
    size: int,
    color: RGBColor,
    *,
    bold: bool = False,
) -> None:
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _paginate(paragraphs: list[str], max_chars: int = 480) -> list[list[str]]:
    pages: list[list[str]] = []
    current: list[str] = []
    used = 0
    for paragraph in paragraphs:
        chunks = textwrap.wrap(
            paragraph,
            width=max_chars,
            break_long_words=True,
            break_on_hyphens=False,
        ) or [""]
        for chunk in chunks:
            if current and used + len(chunk) > max_chars:
                pages.append(current)
                current = []
                used = 0
            current.append(chunk)
            used += len(chunk)
    if current or not pages:
        pages.append(current)
    return pages


def _title_size(title: str, preferred: int) -> int:
    if len(title) <= 45:
        return preferred
    if len(title) <= 80:
        return max(22, preferred - 6)
    return max(17, preferred - 12)
