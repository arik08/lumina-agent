from __future__ import annotations

import hashlib
from collections.abc import Iterator
from contextlib import contextmanager, suppress
from html.parser import HTMLParser
from io import BytesIO
import math
from pathlib import Path, PurePosixPath
from typing import Any
from urllib.parse import urlparse
from zipfile import BadZipFile, ZipFile

from defusedxml import DefusedXmlException
from defusedxml import ElementTree as DefusedElementTree
from sqlalchemy import delete, func, select, update
from sqlalchemy.exc import IntegrityError
from sqlalchemy.orm import Session

from ..api.errors import ApiProblem
from ..document_limits import MAX_DOCUMENT_PAGES, MAX_OPENXML_MEMBERS
from ..authorization import require_conversation, require_project
from ..models import Artifact, ArtifactDraft, ArtifactVersion, User, new_uuid, utc_now
from ..storage import ManagedLocalStorage, StorageError
from .render_validation import ArtifactRenderBackend, verify_artifact_render


TEXT_EDITABLE_ARTIFACT_KINDS = frozenset({"html", "markdown", "text", "json", "csv"})


def require_artifact(
    db: Session, user: User, artifact_id: str, *, write: bool = False
) -> Artifact:
    artifact = db.get(Artifact, artifact_id)
    if artifact is None or artifact.deleted_at is not None:
        raise ApiProblem(404, "not_found", "Artifact를 찾을 수 없습니다.")
    require_project(db, user, artifact.project_id, write=write)
    if artifact.conversation_id:
        require_conversation(db, user, artifact.conversation_id, write=write)
    return artifact


def artifact_summary(
    artifact: Artifact, current: ArtifactVersion | None
) -> dict[str, object]:
    return {
        "id": artifact.id,
        "projectId": artifact.project_id,
        "conversationId": artifact.conversation_id,
        "displayName": artifact.display_name,
        "kind": artifact.kind,
        "mimeType": artifact.mime_type,
        "currentVersion": artifact.current_version_number or 0,
        "validationStatus": current.validation_status if current else "pending",
        "size": current.size_bytes if current else 0,
        "createdAt": artifact.created_at,
        "updatedAt": artifact.updated_at,
    }


def current_artifact_version(db: Session, artifact: Artifact) -> ArtifactVersion | None:
    if artifact.current_version_number is None:
        return None
    return db.scalar(
        select(ArtifactVersion).where(
            ArtifactVersion.artifact_id == artifact.id,
            ArtifactVersion.version_number == artifact.current_version_number,
        )
    )


def create_artifact(
    db: Session,
    storage: ManagedLocalStorage,
    *,
    user: User,
    project_id: str,
    conversation_id: str | None,
    source_run_id: str | None,
    display_name: str,
    kind: str,
    mime_type: str,
    content: bytes,
    change_type: str = "generated",
    change_summary: str = "",
    validation_status: str | None = None,
    renderer_manifest: dict[str, Any] | None = None,
    asset_manifest: list[dict[str, Any]] | None = None,
) -> tuple[Artifact, ArtifactVersion]:
    require_project(db, user, project_id, write=True)
    if conversation_id:
        require_conversation(db, user, conversation_id, write=True)
    artifact = Artifact(
        organization_id=user.organization_id,
        project_id=project_id,
        conversation_id=conversation_id,
        source_run_id=source_run_id,
        created_by_user_id=user.id,
        display_name=display_name,
        kind=kind,
        mime_type=mime_type,
        visibility="private",
    )
    db.add(artifact)
    db.flush()
    computed_status, validation = validate_artifact_content(
        kind=kind, mime_type=mime_type, content=content
    )
    version = _write_version(
        db,
        storage,
        artifact=artifact,
        user=user,
        version_number=1,
        content=content,
        parent=None,
        source=None,
        change_type=change_type,
        change_summary=change_summary,
        validation_status=validation_status or computed_status,
        validation_json=validation,
        renderer_manifest=renderer_manifest,
        asset_manifest=asset_manifest,
    )
    artifact.current_version_number = 1
    try:
        db.flush()
    except BaseException:
        discard_artifact_storage(storage, version.storage_key)
        raise
    return artifact, version


def create_artifact_version(
    db: Session,
    storage: ManagedLocalStorage,
    *,
    user: User,
    artifact_id: str,
    base_version: int,
    content: bytes,
    change_type: str,
    change_summary: str,
    source_version: ArtifactVersion | None = None,
) -> ArtifactVersion:
    artifact = require_artifact(db, user, artifact_id, write=True)
    ensure_artifact_text_editable(artifact)
    if artifact.current_version_number != base_version:
        raise ApiProblem(
            409,
            "artifact_version_conflict",
            "Artifact가 다른 곳에서 변경되었습니다. 최신 버전을 확인해 주세요.",
            details={"currentVersion": artifact.current_version_number},
        )
    parent = db.scalar(
        select(ArtifactVersion).where(
            ArtifactVersion.artifact_id == artifact.id,
            ArtifactVersion.version_number == base_version,
        )
    )
    if parent is None:
        raise ApiProblem(
            404, "version_not_found", "기준 Artifact 버전을 찾을 수 없습니다."
        )
    if change_type == "restore":
        if source_version is None or source_version.artifact_id != artifact.id:
            raise ApiProblem(
                400,
                "artifact_restore_source_invalid",
                "복원할 Artifact 원본 버전이 올바르지 않습니다.",
            )
    elif source_version is not None:
        raise ApiProblem(
            400,
            "artifact_source_version_forbidden",
            "복원 작업이 아닌 버전에는 원본 버전을 지정할 수 없습니다.",
        )
    latest = (
        db.scalar(
            select(func.max(ArtifactVersion.version_number)).where(
                ArtifactVersion.artifact_id == artifact.id
            )
        )
        or 0
    )
    validation_status, validation = validate_artifact_content(
        kind=artifact.kind, mime_type=artifact.mime_type, content=content
    )
    version = _write_version(
        db,
        storage,
        artifact=artifact,
        user=user,
        version_number=latest + 1,
        content=content,
        parent=parent,
        source=source_version,
        change_type=change_type,
        change_summary=change_summary,
        validation_status=validation_status,
        validation_json=validation,
    )
    _set_current_version(
        db,
        storage,
        artifact=artifact,
        base_version=base_version,
        version=version,
    )
    return version


def create_binary_artifact_version(
    db: Session,
    storage: ManagedLocalStorage,
    *,
    user: User,
    artifact_id: str,
    base_version: int,
    mime_type: str,
    content: bytes,
    change_type: str,
    change_summary: str,
    renderer_manifest: dict[str, Any] | None = None,
    asset_manifest: list[dict[str, Any]] | None = None,
) -> tuple[Artifact, ArtifactVersion]:
    """Append immutable bytes to an existing image Artifact with CAS semantics."""
    artifact = require_artifact(db, user, artifact_id, write=True)
    if not artifact.mime_type.startswith("image/"):
        raise ApiProblem(
            409,
            "artifact_image_required",
            "대상 Artifact가 이미지 형식이 아닙니다.",
        )
    if artifact.mime_type != mime_type:
        raise ApiProblem(
            409,
            "artifact_mime_conflict",
            "대상 Artifact와 생성 이미지 형식이 다릅니다.",
        )
    if artifact.current_version_number != base_version:
        raise ApiProblem(
            409,
            "artifact_version_conflict",
            "Artifact가 다른 곳에서 변경되었습니다. 최신 버전을 확인해 주세요.",
            details={"currentVersion": artifact.current_version_number},
        )
    parent = db.scalar(
        select(ArtifactVersion).where(
            ArtifactVersion.artifact_id == artifact.id,
            ArtifactVersion.version_number == base_version,
        )
    )
    if parent is None:
        raise ApiProblem(
            404, "version_not_found", "기준 Artifact 버전을 찾을 수 없습니다."
        )
    latest = (
        db.scalar(
            select(func.max(ArtifactVersion.version_number)).where(
                ArtifactVersion.artifact_id == artifact.id
            )
        )
        or 0
    )
    validation_status, validation = validate_artifact_content(
        kind=artifact.kind, mime_type=mime_type, content=content
    )
    version = _write_version(
        db,
        storage,
        artifact=artifact,
        user=user,
        version_number=latest + 1,
        content=content,
        parent=parent,
        source=None,
        change_type=change_type,
        change_summary=change_summary,
        validation_status=validation_status,
        validation_json=validation,
        renderer_manifest=renderer_manifest,
        asset_manifest=asset_manifest,
    )
    _set_current_version(
        db,
        storage,
        artifact=artifact,
        base_version=base_version,
        version=version,
    )
    return artifact, version


def _write_version(
    db: Session,
    storage: ManagedLocalStorage,
    *,
    artifact: Artifact,
    user: User,
    version_number: int,
    content: bytes,
    parent: ArtifactVersion | None,
    source: ArtifactVersion | None,
    change_type: str,
    change_summary: str,
    validation_status: str,
    validation_json: dict[str, Any],
    renderer_manifest: dict[str, Any] | None = None,
    asset_manifest: list[dict[str, Any]] | None = None,
) -> ArtifactVersion:
    digest = hashlib.sha256(content).hexdigest()
    extension = _safe_extension(artifact.display_name, artifact.mime_type)
    version_id = new_uuid()
    key = (
        f"artifacts/{artifact.id}/v{version_number}/"
        f"{version_id}-{digest}.{extension}"
    )
    stored = storage.put_bytes(key, content, expected_sha256=digest)
    version = ArtifactVersion(
        id=version_id,
        artifact_id=artifact.id,
        version_number=version_number,
        storage_backend="local",
        storage_key=stored.key,
        content_hash=stored.sha256,
        size_bytes=stored.size,
        parent_version_id=parent.id if parent else None,
        source_version_id=source.id if source else None,
        change_type=change_type,
        change_prompt_summary=change_summary or None,
        renderer_manifest_json=renderer_manifest
        or {"renderer": artifact.kind, "version": "1"},
        asset_manifest_json=asset_manifest or [],
        validation_status=validation_status,
        validation_json=validation_json,
        created_by_user_id=user.id,
    )
    db.add(version)
    try:
        db.flush()
    except IntegrityError as exc:
        discard_artifact_storage(storage, stored.key)
        raise ApiProblem(
            409,
            "artifact_version_conflict",
            "Artifact version이 다른 작업에서 먼저 저장되었습니다.",
        ) from exc
    except BaseException:
        discard_artifact_storage(storage, stored.key)
        raise
    return version


def _set_current_version(
    db: Session,
    storage: ManagedLocalStorage,
    *,
    artifact: Artifact,
    base_version: int,
    version: ArtifactVersion,
) -> None:
    try:
        result = db.execute(
            update(Artifact)
            .where(
                Artifact.id == artifact.id,
                Artifact.current_version_number == base_version,
                Artifact.deleted_at.is_(None),
            )
            .values(current_version_number=version.version_number)
            .execution_options(synchronize_session=False)
        )
    except BaseException:
        discard_artifact_storage(storage, version.storage_key)
        raise
    if getattr(result, "rowcount", 0) != 1:
        discard_artifact_storage(storage, version.storage_key)
        db.expire(artifact)
        db.refresh(artifact)
        raise ApiProblem(
            409,
            "artifact_version_conflict",
            "Artifact가 다른 곳에서 변경되었습니다. 최신 버전을 확인해 주세요.",
            details={"currentVersion": artifact.current_version_number},
        )
    db.expire(artifact)
    db.refresh(artifact)


def discard_artifact_storage(storage: ManagedLocalStorage, key: str) -> None:
    with suppress(StorageError):
        storage.delete(key)


@contextmanager
def cleanup_artifact_storage_on_error(
    storage: ManagedLocalStorage,
) -> Iterator[list[str]]:
    storage_keys: list[str] = []
    try:
        yield storage_keys
    except BaseException:
        for storage_key in storage_keys:
            discard_artifact_storage(storage, storage_key)
        raise


def read_artifact_version(
    db: Session,
    storage: ManagedLocalStorage,
    *,
    user: User,
    artifact_id: str,
    version_number: int,
) -> tuple[Artifact, ArtifactVersion, bytes]:
    artifact = require_artifact(db, user, artifact_id)
    version = db.scalar(
        select(ArtifactVersion).where(
            ArtifactVersion.artifact_id == artifact.id,
            ArtifactVersion.version_number == version_number,
        )
    )
    if version is None:
        raise ApiProblem(404, "version_not_found", "Artifact 버전을 찾을 수 없습니다.")
    try:
        content = storage.read_bytes(
            version.storage_key, expected_sha256=version.content_hash
        )
    except StorageError as exc:
        raise ApiProblem(
            503, "artifact_content_missing", "Artifact 원본을 읽을 수 없습니다."
        ) from exc
    return artifact, version, content


def save_draft(
    db: Session,
    storage: ManagedLocalStorage,
    *,
    user: User,
    artifact_id: str,
    base_version: int,
    content: bytes,
    expected_etag: str | None,
) -> tuple[ArtifactDraft, str | None]:
    artifact = require_artifact(db, user, artifact_id, write=True)
    ensure_artifact_text_editable(artifact)
    if artifact.current_version_number != base_version:
        raise ApiProblem(
            409,
            "artifact_version_conflict",
            "기준 버전이 최신 버전과 다릅니다.",
            details={"currentVersion": artifact.current_version_number},
        )
    draft = db.scalar(
        select(ArtifactDraft).where(
            ArtifactDraft.artifact_id == artifact.id,
            ArtifactDraft.user_id == user.id,
        )
    )
    if draft is not None:
        if draft.base_version_number != artifact.current_version_number:
            raise ApiProblem(
                409,
                "artifact_draft_stale",
                "이 편집 초안은 이전 Artifact 버전을 기준으로 합니다.",
                details={
                    "baseVersion": draft.base_version_number,
                    "currentVersion": artifact.current_version_number,
                },
            )
        if expected_etag is None:
            raise ApiProblem(
                428,
                "draft_if_match_required",
                "기존 편집 초안을 변경하려면 If-Match가 필요합니다.",
            )
        if draft.etag != expected_etag:
            raise ApiProblem(
                409,
                "draft_conflict",
                "Artifact 편집 초안이 다른 곳에서 변경되었습니다.",
            )
    elif expected_etag is not None:
        raise ApiProblem(
            409,
            "draft_conflict",
            "지정한 Artifact 편집 초안을 찾을 수 없습니다.",
        )
    previous_storage_key = draft.storage_key if draft is not None else None
    digest = hashlib.sha256(content).hexdigest()
    etag = hashlib.sha256(
        f"{artifact.id}:{user.id}:{base_version}:{digest}".encode()
    ).hexdigest()
    extension = _safe_extension(artifact.display_name, artifact.mime_type)
    write_id = new_uuid()
    key = (
        f"artifact-drafts/{artifact.id}/{user.id}/"
        f"{write_id}-{etag}.{extension}"
    )
    stored = storage.put_bytes(key, content, expected_sha256=digest)
    try:
        if draft is None:
            draft = ArtifactDraft(
                artifact_id=artifact.id,
                user_id=user.id,
                base_version_number=base_version,
                storage_key=stored.key,
                content_hash=stored.sha256,
                etag=etag,
            )
            db.add(draft)
            db.flush()
        else:
            result = db.execute(
                update(ArtifactDraft)
                .where(
                    ArtifactDraft.id == draft.id,
                    ArtifactDraft.etag == expected_etag,
                )
                .values(
                    base_version_number=base_version,
                    storage_key=stored.key,
                    content_hash=stored.sha256,
                    etag=etag,
                    updated_at=utc_now(),
                )
                .execution_options(synchronize_session=False)
            )
            if getattr(result, "rowcount", 0) != 1:
                raise ApiProblem(
                    409,
                    "draft_conflict",
                    "Artifact 편집 초안이 다른 곳에서 변경되었습니다.",
                )
            db.refresh(draft)
    except BaseException:
        discard_artifact_storage(storage, stored.key)
        raise
    return draft, previous_storage_key


def read_user_draft(
    db: Session,
    storage: ManagedLocalStorage,
    *,
    user: User,
    artifact_id: str,
) -> tuple[Artifact, ArtifactDraft, bytes]:
    artifact = require_artifact(db, user, artifact_id)
    ensure_artifact_text_editable(artifact)
    draft = db.scalar(
        select(ArtifactDraft).where(
            ArtifactDraft.artifact_id == artifact.id,
            ArtifactDraft.user_id == user.id,
        )
    )
    if draft is None:
        raise ApiProblem(
            404,
            "artifact_draft_not_found",
            "저장된 Artifact 편집 초안이 없습니다.",
        )
    try:
        content = storage.read_bytes(
            draft.storage_key, expected_sha256=draft.content_hash
        )
    except StorageError as exc:
        raise ApiProblem(
            503,
            "artifact_draft_content_missing",
            "Artifact 편집 초안의 원본을 읽을 수 없습니다.",
        ) from exc
    return artifact, draft, content


def delete_user_draft_if_matches(
    db: Session,
    *,
    user: User,
    artifact_id: str,
    base_version: int,
    etag: str | None,
    content_hash: str,
) -> str | None:
    if etag is None:
        return None
    storage_key = db.scalar(
        select(ArtifactDraft.storage_key).where(
            ArtifactDraft.artifact_id == artifact_id,
            ArtifactDraft.user_id == user.id,
            ArtifactDraft.base_version_number == base_version,
            ArtifactDraft.etag == etag,
            ArtifactDraft.content_hash == content_hash,
        )
    )
    if storage_key is None:
        return None
    result = db.execute(
        delete(ArtifactDraft).where(
            ArtifactDraft.artifact_id == artifact_id,
            ArtifactDraft.user_id == user.id,
            ArtifactDraft.base_version_number == base_version,
            ArtifactDraft.etag == etag,
            ArtifactDraft.content_hash == content_hash,
        )
    )
    return storage_key if getattr(result, "rowcount", 0) == 1 else None


_MAX_OPENXML_UNCOMPRESSED_BYTES = 512 * 1024 * 1024
_MAX_OPENXML_ENTRY_BYTES = 128 * 1024 * 1024
_MAX_OPENXML_RELATIONSHIP_BYTES = 2 * 1024 * 1024
_MAX_XLSX_VISITED_CELLS = 2_000_000
_MIN_PDF_PAGE_POINTS = 36.0
_MAX_PDF_PAGE_POINTS = 7_500.0
_MAX_PDF_PAGE_AREA_POINTS = 22_500_000.0
_OPENXML_REQUIRED_PARTS = {
    "docx": "word/document.xml",
    "xlsx": "xl/workbook.xml",
    "pptx": "ppt/presentation.xml",
}


def _inspect_openxml_package(
    content: bytes, *, kind: str
) -> tuple[list[str], dict[str, int]]:
    errors: list[str] = []
    details = {"packageParts": 0, "externalHyperlinks": 0}
    try:
        with ZipFile(BytesIO(content)) as package:
            entries = package.infolist()
            details["packageParts"] = len(entries)
            if len(entries) > MAX_OPENXML_MEMBERS:
                errors.append("openxml_package_file_limit_exceeded")
            names: set[str] = set()
            total_uncompressed = 0
            for entry in entries:
                normalized = entry.filename.replace("\\", "/")
                path = PurePosixPath(normalized)
                if (
                    normalized.startswith("/")
                    or ".." in path.parts
                    or (path.parts and ":" in path.parts[0])
                ):
                    errors.append("unsafe_openxml_path")
                lowered = normalized.casefold()
                if lowered in names:
                    errors.append("openxml_duplicate_part")
                names.add(lowered)
                total_uncompressed += entry.file_size
                if entry.file_size > _MAX_OPENXML_ENTRY_BYTES:
                    errors.append("openxml_package_entry_too_large")
                if entry.flag_bits & 0x1:
                    errors.append("openxml_encrypted_part")
                if (
                    entry.file_size > 1024 * 1024
                    and entry.file_size > max(entry.compress_size, 1) * 200
                ):
                    errors.append("openxml_compression_ratio_exceeded")
                if lowered.endswith("vbaproject.bin") or "/macros/" in lowered:
                    errors.append("openxml_macros_forbidden")
            if total_uncompressed > _MAX_OPENXML_UNCOMPRESSED_BYTES:
                errors.append("openxml_package_too_large")
            required = {
                "[content_types].xml",
                "_rels/.rels",
                _OPENXML_REQUIRED_PARTS[kind],
            }
            if not required.issubset(names):
                errors.append(f"invalid_{kind}_structure")

            if errors:
                return list(dict.fromkeys(errors)), details

            for entry in entries:
                if not entry.filename.casefold().endswith(".rels"):
                    continue
                if entry.file_size > _MAX_OPENXML_RELATIONSHIP_BYTES:
                    errors.append("openxml_relationships_too_large")
                    continue
                try:
                    root = DefusedElementTree.fromstring(package.read(entry))
                except (DefusedXmlException, DefusedElementTree.ParseError, OSError):
                    errors.append("invalid_openxml_relationships")
                    continue
                for relationship in root:
                    if (
                        relationship.attrib.get("TargetMode", "").casefold()
                        != "external"
                    ):
                        continue
                    relationship_type = relationship.attrib.get("Type", "")
                    target = relationship.attrib.get("Target", "")
                    if not relationship_type.casefold().endswith("/hyperlink"):
                        errors.append("openxml_external_resource_forbidden")
                        continue
                    details["externalHyperlinks"] += 1
                    if not _is_safe_external_link(target):
                        errors.append("unsafe_external_link")
    except (BadZipFile, KeyError, OSError):
        errors.append(f"invalid_{kind}_structure")
    return list(dict.fromkeys(errors)), details


def _is_safe_external_link(target: str) -> bool:
    parsed = urlparse(target.strip())
    return parsed.scheme.casefold() in {"http", "https", "mailto"}


def _inspect_pdf_links(reader: Any) -> tuple[int, list[str]]:
    link_count = 0
    errors: list[str] = []
    try:
        catalog = reader.trailer["/Root"].get_object()
        if catalog.get("/OpenAction") is not None or catalog.get("/AA") is not None:
            errors.append("unsafe_pdf_action")
        names = catalog.get("/Names")
        if names is not None:
            names = names.get_object()
            if names.get("/JavaScript") is not None:
                errors.append("unsafe_pdf_action")
    except Exception:
        errors.append("invalid_pdf_catalog")
    for page in reader.pages:
        try:
            if page.get("/AA") is not None:
                errors.append("unsafe_pdf_action")
            annotations = page.get("/Annots") or []
            for annotation_ref in annotations:
                annotation = annotation_ref.get_object()
                if str(annotation.get("/Subtype", "")) != "/Link":
                    continue
                link_count += 1
                action = annotation.get("/A")
                if action is None:
                    continue
                if annotation.get("/AA") is not None:
                    errors.append("unsafe_pdf_action")
                action = action.get_object()
                action_type = str(action.get("/S", ""))
                if action_type in {"/JavaScript", "/Launch", "/SubmitForm"}:
                    errors.append("unsafe_pdf_action")
                uri = action.get("/URI")
                if uri is not None and not _is_safe_external_link(str(uri)):
                    errors.append("unsafe_external_link")
        except Exception:
            errors.append("invalid_pdf_annotations")
    return link_count, list(dict.fromkeys(errors))


class _ArtifactHTMLValidator(HTMLParser):
    forbidden_tags = {"iframe", "object", "embed", "base"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.tags: set[str] = set()
        self.errors: list[str] = []
        self._title_depth = 0
        self._html_closed = False
        self.title_text: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        lowered = tag.casefold()
        if self._html_closed:
            self.errors.append("trailing_content_after_html")
        self.tags.add(lowered)
        if lowered in self.forbidden_tags:
            self.errors.append(f"forbidden_tag:{lowered}")
        if lowered == "title":
            self._title_depth += 1
        for name, value in attrs:
            attr_name = name.casefold()
            attr_value = (value or "").strip().casefold()
            if attr_name in {"href", "src", "action", "formaction"} and (
                attr_value.startswith("data:text/html")
            ):
                self.errors.append(f"unsafe_url:{attr_name}")

    def handle_endtag(self, tag: str) -> None:
        lowered = tag.casefold()
        if lowered == "title" and self._title_depth:
            self._title_depth -= 1
        elif lowered == "html":
            self._html_closed = True

    def handle_data(self, data: str) -> None:
        if self._title_depth:
            self.title_text.append(data)
        if self._html_closed and data.strip():
            self.errors.append("trailing_content_after_html")


def validate_artifact_content(
    *,
    kind: str,
    mime_type: str,
    content: bytes,
    render_backend: ArtifactRenderBackend | None = None,
) -> tuple[str, dict[str, Any]]:
    """Validate bytes independently from generation and persist reproducible checks."""
    checks = ["content_hash", "non_empty"]
    errors: list[str] = []
    warnings: list[str] = []
    details: dict[str, Any] = {}
    if not content or not content.strip():
        errors.append("empty_content")

    if mime_type == "text/html" or kind == "html":
        checks.extend(["utf8", "html_structure", "executable_content"])
        try:
            source = content.decode("utf-8", errors="strict")
        except UnicodeDecodeError:
            errors.append("invalid_utf8")
        else:
            validator = _ArtifactHTMLValidator()
            try:
                validator.feed(source)
                validator.close()
            except Exception:
                errors.append("html_parse_error")
            required = {"html", "head", "title", "body"}
            missing = sorted(required - validator.tags)
            errors.extend(f"missing_tag:{tag}" for tag in missing)
            if not "".join(validator.title_text).strip():
                errors.append("empty_title")
            errors.extend(validator.errors)
            details["htmlTitle"] = "".join(validator.title_text).strip()
    elif mime_type.startswith("text/") or kind in {
        "markdown",
        "text",
        "json",
        "csv",
    }:
        checks.append("utf8")
        try:
            content.decode("utf-8", errors="strict")
        except UnicodeDecodeError:
            errors.append("invalid_utf8")
    elif mime_type == "application/pdf" or kind == "pdf":
        checks.extend(
            [
                "pdf_signature",
                "pdf_structure",
                "pdf_text",
                "pdf_page_dimensions",
                "pdf_links",
            ]
        )
        if not content.startswith(b"%PDF-"):
            errors.append("invalid_pdf_signature")
        else:
            from pypdf import PdfReader

            try:
                reader = PdfReader(BytesIO(content))
            except Exception:
                errors.append("invalid_pdf_structure")
            else:
                if reader.is_encrypted:
                    try:
                        decrypted = bool(reader.decrypt(""))
                    except Exception:
                        decrypted = False
                    if not decrypted:
                        errors.append("encrypted_pdf")
                page_count = len(reader.pages)
                if page_count == 0:
                    errors.append("empty_pdf")
                if page_count > MAX_DOCUMENT_PAGES:
                    errors.append("pdf_page_limit_exceeded")
                page_text: list[str] = []
                if not errors:
                    try:
                        for page_number, page in enumerate(reader.pages, start=1):
                            width = abs(float(page.mediabox.width))
                            height = abs(float(page.mediabox.height))
                            if (
                                not math.isfinite(width)
                                or not math.isfinite(height)
                                or width < _MIN_PDF_PAGE_POINTS
                                or height < _MIN_PDF_PAGE_POINTS
                                or width > _MAX_PDF_PAGE_POINTS
                                or height > _MAX_PDF_PAGE_POINTS
                                or width * height > _MAX_PDF_PAGE_AREA_POINTS
                            ):
                                errors.append(
                                    "pdf_page_size_out_of_range:"
                                    f"{page_number}:{width:g}x{height:g}"
                                )
                            page_text.append(page.extract_text() or "")
                    except Exception:
                        errors.append("invalid_pdf_structure")
                extracted_text = "".join(page_text)
                if not extracted_text.strip():
                    errors.append("missing_pdf_text")
                link_count, link_errors = _inspect_pdf_links(reader)
                errors.extend(link_errors)
                details["pageCount"] = page_count
                details["extractedCharacters"] = len(extracted_text)
                details["linkCount"] = link_count
    elif mime_type == "image/png" or kind == "image_png":
        checks.append("png_signature")
        if not content.startswith(b"\x89PNG\r\n\x1a\n"):
            errors.append("invalid_png_signature")
    elif mime_type == "image/jpeg" or kind == "image_jpeg":
        checks.append("jpeg_signature")
        if not content.startswith(b"\xff\xd8\xff"):
            errors.append("invalid_jpeg_signature")
    elif mime_type == "image/webp" or kind == "image_webp":
        checks.append("webp_signature")
        if not (content.startswith(b"RIFF") and content[8:12] == b"WEBP"):
            errors.append("invalid_webp_signature")
    elif (
        mime_type
        == ("application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        or kind == "docx"
    ):
        checks.extend(
            [
                "openxml_zip_signature",
                "openxml_package_structure",
                "openxml_macro_policy",
                "openxml_external_relationships",
            ]
        )
        if not content.startswith(b"PK"):
            errors.append("invalid_openxml_signature")
        else:
            from docx import Document as DocxDocument

            checks.extend(["docx_structure", "docx_content", "docx_page_layout"])
            package_errors, package_details = _inspect_openxml_package(
                content, kind="docx"
            )
            errors.extend(package_errors)
            details.update(package_details)
            if not package_errors:
                try:
                    document = DocxDocument(BytesIO(content))
                except Exception:
                    errors.append("invalid_docx_structure")
                else:
                    text = "".join(paragraph.text for paragraph in document.paragraphs)
                    text += "".join(
                        cell.text
                        for table in document.tables
                        for row in table.rows
                        for cell in row.cells
                    )
                    if not text.strip():
                        errors.append("empty_docx")
                    if not document.sections or any(
                        section.page_width is None
                        or section.page_height is None
                        or section.page_width <= 0
                        or section.page_height <= 0
                        for section in document.sections
                    ):
                        errors.append("invalid_docx_page_layout")
                    details["sectionCount"] = len(document.sections)
                    details["paragraphCount"] = len(document.paragraphs)
                    details["tableCount"] = len(document.tables)
                    details["inlineShapeCount"] = len(document.inline_shapes)
    elif (
        mime_type
        == ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        or kind == "xlsx"
    ):
        checks.extend(
            [
                "openxml_zip_signature",
                "openxml_package_structure",
                "openxml_macro_policy",
                "openxml_external_relationships",
            ]
        )
        if not content.startswith(b"PK"):
            errors.append("invalid_openxml_signature")
        else:
            from openpyxl import load_workbook

            checks.extend(
                [
                    "xlsx_structure",
                    "xlsx_visible_sheet",
                    "xlsx_populated_cells",
                    "xlsx_formula_references",
                ]
            )
            package_errors, package_details = _inspect_openxml_package(
                content, kind="xlsx"
            )
            errors.extend(package_errors)
            details.update(package_details)
            if not package_errors:
                try:
                    workbook = load_workbook(
                        BytesIO(content), data_only=False, keep_links=False
                    )
                except Exception:
                    errors.append("invalid_xlsx_structure")
                else:
                    try:
                        if not workbook.worksheets:
                            errors.append("empty_xlsx")
                        if not any(
                            sheet.sheet_state == "visible"
                            for sheet in workbook.worksheets
                        ):
                            errors.append("missing_visible_xlsx_sheet")
                        populated_cells = 0
                        hyperlink_count = 0
                        chart_count = 0
                        for sheet in workbook.worksheets:
                            visited_cells = sheet.max_row * sheet.max_column
                            if visited_cells > _MAX_XLSX_VISITED_CELLS:
                                errors.append(f"xlsx_sheet_too_large:{sheet.title}")
                                continue
                            chart_count += len(getattr(sheet, "_charts", ()))
                            for row in sheet.iter_rows():
                                for cell in row:
                                    if cell.value is not None:
                                        populated_cells += 1
                                    if cell.data_type == "f":
                                        formula = str(cell.value)
                                        if "#REF!" in formula.upper():
                                            errors.append(
                                                "xlsx_formula_reference:"
                                                f"{sheet.title}:{cell.coordinate}"
                                            )
                                        if "[" in formula and "]" in formula:
                                            errors.append(
                                                "xlsx_external_formula_reference:"
                                                f"{sheet.title}:{cell.coordinate}"
                                            )
                                    if cell.hyperlink is not None:
                                        hyperlink_count += 1
                                        if not _is_safe_external_link(
                                            str(cell.hyperlink.target)
                                        ):
                                            errors.append("unsafe_external_link")
                        if not populated_cells:
                            errors.append("empty_xlsx")
                        details["sheetNames"] = workbook.sheetnames
                        details["sheetCount"] = len(workbook.worksheets)
                        details["populatedCells"] = populated_cells
                        details["hyperlinkCount"] = hyperlink_count
                        details["chartCount"] = chart_count
                    finally:
                        workbook.close()
    elif (
        mime_type
        == ("application/vnd.openxmlformats-officedocument.presentationml.presentation")
        or kind == "pptx"
    ):
        checks.extend(
            [
                "openxml_zip_signature",
                "openxml_package_structure",
                "openxml_macro_policy",
                "openxml_external_relationships",
            ]
        )
        if not content.startswith(b"PK"):
            errors.append("invalid_openxml_signature")
        else:
            from pptx import Presentation

            checks.extend(
                [
                    "pptx_structure",
                    "pptx_content",
                    "pptx_shape_bounds",
                    "pptx_editable_text",
                ]
            )
            package_errors, package_details = _inspect_openxml_package(
                content, kind="pptx"
            )
            errors.extend(package_errors)
            details.update(package_details)
            if not package_errors:
                try:
                    presentation = Presentation(BytesIO(content))
                except Exception:
                    errors.append("invalid_pptx_structure")
                else:
                    if not presentation.slides:
                        errors.append("empty_pptx")
                    if (
                        presentation.slide_width is None
                        or presentation.slide_height is None
                        or presentation.slide_width <= 0
                        or presentation.slide_height <= 0
                    ):
                        errors.append("invalid_pptx_page_layout")
                    editable_text = 0
                    hyperlink_count = 0
                    for slide_number, slide in enumerate(presentation.slides, start=1):
                        for shape in slide.shapes:
                            if (
                                getattr(shape, "has_text_frame", False)
                                and shape.text.strip()
                            ):
                                editable_text += 1
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        address = run.hyperlink.address
                                        if address:
                                            hyperlink_count += 1
                                            if not _is_safe_external_link(address):
                                                errors.append("unsafe_external_link")
                            click_action = getattr(shape, "click_action", None)
                            address = (
                                click_action.hyperlink.address
                                if click_action is not None
                                else None
                            )
                            if address:
                                hyperlink_count += 1
                                if not _is_safe_external_link(address):
                                    errors.append("unsafe_external_link")
                            if (
                                shape.left < 0
                                or shape.top < 0
                                or shape.width <= 0
                                or shape.height <= 0
                                or shape.left + shape.width > presentation.slide_width
                                or shape.top + shape.height > presentation.slide_height
                            ):
                                errors.append(
                                    "pptx_shape_out_of_bounds:"
                                    f"{slide_number}:{shape.shape_id}"
                                )
                    if not editable_text:
                        errors.append("missing_pptx_editable_text")
                    details["slideCount"] = len(presentation.slides)
                    details["editableTextShapes"] = editable_text
                    details["hyperlinkCount"] = hyperlink_count

    unique_errors = list(dict.fromkeys(errors))
    render_verification_required = kind.casefold() in {
        "html",
        "pdf",
        "docx",
        "xlsx",
        "pptx",
    } or mime_type in {
        "text/html",
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    render_verified = False
    renderer: str | None = None
    pages: list[dict[str, object]] = []
    if render_verification_required and not unique_errors:
        if mime_type == "text/html" or kind == "html":
            warnings.append("render_verification_pending")
        else:
            expected_page_count = details.get("pageCount")
            render = verify_artifact_render(
                kind=kind,
                mime_type=mime_type,
                content=content,
                expected_page_count=(
                    int(expected_page_count)
                    if isinstance(expected_page_count, int)
                    else None
                ),
                backend=render_backend,
            )
            render_verified = render.render_verified
            renderer = render.renderer
            pages = render.pages
            warnings.extend(render.warnings)
            unique_errors.extend(render.errors)
            checks.extend(render.checks)
    unique_errors = list(dict.fromkeys(unique_errors))
    unique_warnings = list(dict.fromkeys(warnings))
    return (
        (
            "failed"
            if unique_errors
            else "passed"
            if render_verified or not render_verification_required
            else "structural_passed"
        ),
        {
            "validatorVersion": "lumina-artifact-v3",
            "verificationLevel": "render" if render_verified else "structural",
            "renderVerified": render_verified,
            "renderVerificationRequired": render_verification_required,
            "renderer": renderer,
            "pages": pages,
            "warnings": unique_warnings,
            "checks": checks,
            "errors": unique_errors,
            "contentHash": hashlib.sha256(content).hexdigest(),
            "sizeBytes": len(content),
            "details": details,
        },
    )


def ensure_artifact_text_editable(artifact: Artifact) -> None:
    if artifact.kind.casefold() not in TEXT_EDITABLE_ARTIFACT_KINDS:
        raise ApiProblem(
            409,
            "artifact_binary_edit_unsupported",
            "바이너리 Artifact는 sourceText 방식으로 편집할 수 없습니다.",
        )


def _safe_extension(display_name: str, mime_type: str) -> str:
    suffix = Path(display_name).suffix.removeprefix(".").lower()
    if suffix.isalnum() and 1 <= len(suffix) <= 12:
        return suffix
    return {
        "text/html": "html",
        "text/markdown": "md",
        "application/json": "json",
        "text/plain": "txt",
        "text/csv": "csv",
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/webp": "webp",
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    }.get(mime_type, "bin")
