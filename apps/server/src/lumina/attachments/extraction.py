from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from itertools import islice
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PasswordType, PdfReader

from ..document_limits import MAX_DOCUMENT_PAGES


MAX_EXTRACTED_CHARS = 2_000_000
MAX_SPREADSHEET_CELLS = 50_000


@dataclass(frozen=True, slots=True)
class ExtractionResult:
    status: str
    text: str = ""
    locator_map: dict[str, Any] = field(default_factory=dict)
    metadata: dict[str, Any] = field(default_factory=dict)


def extract_attachment_text(
    *, filename: str, mime_type: str, content: bytes
) -> ExtractionResult:
    try:
        if mime_type.startswith("text/"):
            text = content.decode("utf-8")
            return _completed(
                text,
                locator_map={"kind": "line", "count": len(text.splitlines())},
            )
        if mime_type == "application/pdf":
            return extract_pdf_text(content=content)
        if mime_type.endswith("wordprocessingml.document"):
            return _extract_docx(content)
        if mime_type.endswith("spreadsheetml.sheet"):
            return _extract_xlsx(content)
        if mime_type.endswith("presentationml.presentation"):
            return _extract_pptx(content)
        if mime_type.startswith("image/"):
            return ExtractionResult(status="not_required")
        return ExtractionResult(
            status="unsupported",
            metadata={"reason": f"unsupported:{Path(filename).suffix.lower()}"},
        )
    except Exception as exc:
        return ExtractionResult(
            status="failed",
            metadata={"errorType": type(exc).__name__},
        )


def extract_pdf_text(
    *,
    content: bytes,
    page_start: int = 1,
    page_end: int | None = None,
) -> ExtractionResult:
    try:
        if page_start < 1 or (page_end is not None and page_end < page_start):
            raise ValueError("invalid PDF page range")
        reader = PdfReader(BytesIO(content))
        if (
            reader.is_encrypted
            and reader.decrypt("") == PasswordType.NOT_DECRYPTED
        ):
            return ExtractionResult(
                status="failed",
                metadata={"errorType": "EncryptedPdf"},
            )
        total_pages = len(reader.pages)
        if page_start > total_pages:
            raise ValueError("PDF page range starts after the final page")
        selected_end = min(
            page_end or total_pages,
            total_pages,
            page_start + MAX_DOCUMENT_PAGES - 1,
        )
        pages = [
            f"[Page {index}]\n{reader.pages[index - 1].extract_text() or ''}"
            for index in range(page_start, selected_end + 1)
        ]
        return _completed(
            "\n\n".join(pages),
            locator_map={
                "kind": "page",
                "count": total_pages,
                "start": page_start,
                "end": selected_end,
            },
            metadata={
                "hasMorePages": selected_end < total_pages,
                "truncatedByPageLimit": (
                    selected_end < min(page_end or total_pages, total_pages)
                ),
            },
        )
    except Exception as exc:
        return ExtractionResult(
            status="failed",
            metadata={"errorType": type(exc).__name__},
        )


def _extract_docx(content: bytes) -> ExtractionResult:
    document = Document(BytesIO(content))
    blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table_index, table in enumerate(document.tables, start=1):
        blocks.append(f"[Table {table_index}]")
        for row in table.rows:
            blocks.append("\t".join(cell.text for cell in row.cells))
    return _completed(
        "\n".join(blocks),
        locator_map={
            "kind": "block",
            "paragraphCount": len(document.paragraphs),
            "tableCount": len(document.tables),
        },
    )


def _extract_xlsx(content: bytes) -> ExtractionResult:
    workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
    blocks: list[str] = []
    cell_count = 0
    truncated = False
    try:
        for sheet in workbook.worksheets:
            blocks.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                cell_count += len(values)
                blocks.append("\t".join(values))
                if cell_count >= MAX_SPREADSHEET_CELLS:
                    truncated = True
                    break
            if truncated:
                break
    finally:
        workbook.close()
    return _completed(
        "\n".join(blocks),
        locator_map={"kind": "sheet", "count": len(workbook.sheetnames)},
        metadata={"cellCount": cell_count, "truncatedByCellLimit": truncated},
    )


def _extract_pptx(content: bytes) -> ExtractionResult:
    presentation = Presentation(BytesIO(content))
    slides: list[str] = []
    total_slides = len(presentation.slides)
    for index, slide in enumerate(
        islice(presentation.slides, MAX_DOCUMENT_PAGES), start=1
    ):
        text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        slides.append(f"[Slide {index}]\n" + "\n".join(text))
    return _completed(
        "\n\n".join(slides),
        locator_map={"kind": "slide", "count": total_slides},
        metadata={"truncatedBySlideLimit": total_slides > MAX_DOCUMENT_PAGES},
    )


def _completed(
    text: str,
    *,
    locator_map: dict[str, Any],
    metadata: dict[str, Any] | None = None,
) -> ExtractionResult:
    normalized = text.replace("\x00", "").strip()
    truncated = len(normalized) > MAX_EXTRACTED_CHARS
    if truncated:
        normalized = normalized[:MAX_EXTRACTED_CHARS]
    return ExtractionResult(
        status="completed",
        text=normalized,
        locator_map=locator_map,
        metadata={**(metadata or {}), "truncated": truncated},
    )
