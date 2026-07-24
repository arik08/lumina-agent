from __future__ import annotations

import time
from io import BytesIO
from pathlib import Path

import pytest
from docx import Document
from fastapi.testclient import TestClient
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy.orm import Session

from lumina.agent.executor import (
    _REPORT_TOOL_SCHEMA,
    _append_report_fragment,
    local_run_executor,
)
from lumina.agent.tool_schemas import (
    _EXTEND_REPORT_TOOL_SCHEMA,
    _report_tool_schema,
)
from lumina.artifacts.render_validation import LocalArtifactRenderBackend
from lumina.artifacts.reporting import REPORT_FORMATS, generate_report
from lumina.artifacts.service import validate_artifact_content
from lumina.config import Settings
from lumina.main import create_app
from lumina.models import ArtifactVersion
from lumina.providers import MockProvider, MockToolCall


_EXPECTED_METADATA = {
    "html": ("광양_설비_점검_보고서.html", "html", "text/html"),
    "markdown": ("광양_설비_점검_보고서.md", "markdown", "text/markdown"),
    "docx": (
        "광양_설비_점검_보고서.docx",
        "docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ),
    "xlsx": (
        "광양_설비_점검_보고서.xlsx",
        "xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ),
    "pptx": (
        "광양_설비_점검_보고서.pptx",
        "pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ),
    "pdf": ("광양_설비_점검_보고서.pdf", "pdf", "application/pdf"),
}


@pytest.fixture(autouse=True)
def _disable_host_renderers(monkeypatch: pytest.MonkeyPatch) -> None:
    """Keep unit/API expectations deterministic; installed tools have separate tests."""
    monkeypatch.setattr(
        LocalArtifactRenderBackend,
        "find_executable",
        lambda _self, _candidates: None,
    )


def _arguments(report_format: str) -> dict[str, object]:
    arguments: dict[str, object] = {
        "format": report_format,
        "title": "광양 설비 점검 보고서",
        "executive_summary": "핵심 설비 2건을 확인했습니다.",
        "key_metrics": [{"label": "점검 설비", "value": "2건"}],
        "sections": [
            {
                "heading": "점검 결과",
                "body": "이상 징후 1건을 확인했습니다.",
                "bullets": ["베어링 온도를 확인합니다.", "담당자가 재점검합니다."],
            }
        ],
        "action_items": ["48시간 안에 재점검합니다."],
    }
    if report_format == "html":
        arguments["html_source"] = """<!doctype html>
<html lang="ko"><head><meta charset="utf-8"><title>광양 설비 점검 보고서</title>
<style>@media print{body{color:#000}}</style></head>
<body><main><h1>광양 설비 점검 보고서</h1>
<section><h2>점검 결과</h2><p>이상 징후 1건을 확인했습니다.</p></section>
</main></body></html>"""
    return arguments


@pytest.mark.parametrize("report_format", REPORT_FORMATS)
def test_report_formats_reopen_and_retain_korean_content(report_format: str) -> None:
    report = generate_report("설비 상태를 정리해 주세요.", _arguments(report_format))
    expected_name, expected_kind, expected_mime = _EXPECTED_METADATA[report_format]

    assert report.display_name == expected_name
    assert report.kind == expected_kind
    assert report.mime_type == expected_mime
    _assert_reopened(report_format, report.content)

    status, validation = validate_artifact_content(
        kind=report.kind,
        mime_type=report.mime_type,
        content=report.content,
    )
    expected_status = "passed" if report_format == "markdown" else "structural_passed"
    assert status == expected_status, validation
    assert validation["errors"] == []
    assert validation["renderVerified"] is False
    assert validation["renderVerificationRequired"] is (report_format != "markdown")
    detail_keys = {
        "html": "htmlTitle",
        "docx": "paragraphCount",
        "xlsx": "populatedCells",
        "pptx": "slideCount",
        "pdf": "pageCount",
    }
    if report_format in detail_keys:
        assert validation["details"][detail_keys[report_format]]


def test_create_report_schema_advertises_every_supported_format() -> None:
    report_description = _REPORT_TOOL_SCHEMA["function"]["description"]
    schema = _REPORT_TOOL_SCHEMA["function"]["parameters"]["properties"]["format"]
    assert schema["enum"] == list(REPORT_FORMATS)
    title_description = _REPORT_TOOL_SCHEMA["function"]["parameters"]["properties"][
        "title"
    ]["description"]
    assert "Artifact filename" in title_description
    assert "generic names" in title_description
    html_description = _REPORT_TOOL_SCHEMA["function"]["parameters"]["properties"][
        "html_source"
    ]["description"]
    assert "`.mermaid` element" in html_description
    assert "expand/zoom controls" in html_description
    assert "do not include a Mermaid CDN script" in html_description
    assert "one hue family per top-level semantic branch" in html_description
    assert "Assign a class to every node" in html_description
    assert "at least 4.5:1 contrast" in html_description
    for color in (
        "#3288bd",
        "#66c2a5",
        "#e6f598",
        "#d53e4f",
        "#9e0142",
        "#f46d43",
        "#fdae61",
        "#fee08b",
        "#abdda4",
        "#5e4fa2",
    ):
        assert color in report_description
        assert color in html_description
    assert "user's designated default visual palette" in report_description
    assert "substituting Lumina app cobalt or an all-gray scheme" in html_description


def test_create_report_schema_separates_html_from_structured_report_fields() -> None:
    parameters = _REPORT_TOOL_SCHEMA["function"]["parameters"]

    assert parameters["required"] == ["format", "title"]
    assert parameters["oneOf"] == [
        {
            "properties": {"format": {"const": "html"}},
            "required": ["html_source"],
        },
        {
            "properties": {
                "format": {
                    "enum": [
                        report_format
                        for report_format in REPORT_FORMATS
                        if report_format != "html"
                    ]
                }
            },
            "required": ["executive_summary", "sections"],
        },
    ]
    assert "legacy executive_summary" in parameters["description"]
    assert "action_items is optional" in parameters["description"]


def test_large_html_report_schema_keeps_html_source_required_without_legacy_fields() -> (
    None
):
    parameters = _report_tool_schema(30_000)["function"]["parameters"]

    assert parameters["properties"]["html_source"]["minLength"] == 48_000
    assert parameters["oneOf"][0]["required"] == ["html_source"]
    assert "sections" not in parameters["oneOf"][0]["required"]


def test_extend_report_appends_fragments_without_accepting_full_html() -> None:
    html = (
        "<!doctype html><html><head><title>원본</title></head>"
        "<body><main><h1>원본</h1></main></body></html>"
    )
    fragment = "<section id='added'><h2>추가 분석</h2></section>"

    extended = _append_report_fragment(html, fragment, mime_type="text/html")

    assert extended.startswith(html.split("</main>", 1)[0])
    assert f"{fragment}\n</main>" in extended
    assert extended.count("<!doctype html>") == 1
    with pytest.raises(ValueError, match="body fragments only"):
        _append_report_fragment(
            html,
            "<html><body><p>전체 문서 재생성</p></body></html>",
            mime_type="text/html",
        )
    assert _append_report_fragment(
        "# 원본\n",
        "## 추가 분석",
        mime_type="text/markdown",
    ) == "# 원본\n\n## 추가 분석\n"
    assert _EXTEND_REPORT_TOOL_SCHEMA["function"]["parameters"]["required"] == [
        "content"
    ]


def test_html_source_preserves_visual_artifact_and_executable_javascript() -> None:
    source = """<!doctype html>
<html lang="ko"><head><meta charset="utf-8"><title>시각 보고서</title>
<style>body{background:#eef3fa}.chart{display:grid}</style></head>
<body><main><section class="chart"><h1>시각 보고서</h1><svg viewBox="0 0 10 10" aria-label="추세"><path d="M0 9L9 1"/></svg></section></main></body></html>"""
    arguments = _arguments("html")
    arguments["html_source"] = source

    report = generate_report("시각 보고서를 만들어 주세요.", arguments)

    assert report.content.decode("utf-8") == source
    assert report.display_name == "광양_설비_점검_보고서.html"
    assert "html_source" in _REPORT_TOOL_SCHEMA["function"]["parameters"]["properties"]

    executable_source = source.replace(
        "</body>",
        "<script>document.body.dataset.ready = 'true';</script></body>",
    )
    arguments["html_source"] = executable_source

    executable_report = generate_report("시각 보고서를 만들어 주세요.", arguments)

    assert executable_report.content.decode("utf-8") == executable_source


def test_html_report_without_html_source_is_rejected_instead_of_flattened() -> None:
    arguments = _arguments("html")
    arguments.pop("html_source")

    with pytest.raises(ValueError, match="html_source로 제공"):
        generate_report("시각 보고서를 만들어 주세요.", arguments)


def test_non_html_report_omits_follow_up_section_when_action_items_are_absent() -> None:
    arguments = _arguments("markdown")
    arguments.pop("action_items")
    arguments.pop("key_metrics")

    report = generate_report("시장 동향을 분석해 주세요.", arguments)
    source = report.content.decode("utf-8")

    assert "후속 조치" not in source
    assert "검토 섹션" in source
    assert "문서 형식 검증" in source


def test_report_filename_is_safe_and_does_not_repeat_the_extension() -> None:
    arguments = _arguments("html")
    arguments["title"] = "  광양 / 포항: 설비 비교.html  "

    report = generate_report("두 사업장을 비교해 주세요.", arguments)

    assert report.display_name == "광양_포항_설비_비교.html"


def test_html_source_is_rejected_for_non_html_reports() -> None:
    arguments = _arguments("pdf")
    arguments["html_source"] = (
        "<!doctype html><html><head><title>x</title></head><body>x</body></html>"
    )
    with pytest.raises(ValueError, match="HTML 보고서에서만"):
        generate_report("PDF 보고서를 만들어 주세요.", arguments)


def test_xlsx_keeps_formula_like_model_text_as_plain_text() -> None:
    arguments = _arguments("xlsx")
    arguments["sections"] = [
        {
            "heading": "검토 결과",
            "body": '=HYPERLINK("https://example.com", "열기")',
            "bullets": ["=1+1"],
        }
    ]
    report = generate_report("수식처럼 보이는 텍스트를 정리해 주세요.", arguments)
    workbook = load_workbook(BytesIO(report.content), data_only=False)
    formula_cells = [
        cell.coordinate
        for sheet in workbook.worksheets
        for row in sheet.iter_rows()
        for cell in row
        if cell.data_type == "f"
    ]
    values = [
        cell.value
        for sheet in workbook.worksheets
        for row in sheet.iter_rows()
        for cell in row
        if cell.value is not None
    ]
    workbook.close()

    assert formula_cells == []
    assert '=HYPERLINK("https://example.com", "열기")' in values
    assert "• =1+1" in values


def test_unknown_report_format_is_rejected() -> None:
    with pytest.raises(ValueError, match="지원하지 않는 보고서 형식"):
        generate_report("요청", _arguments("rtf"))


@pytest.mark.parametrize("report_format", REPORT_FORMATS)
def test_create_report_tool_persists_and_downloads_selected_format(
    report_format: str, monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = Settings(
        environment="test",
        database_url=f"sqlite:///{(tmp_path / f'{report_format}.db').as_posix()}",
        data_dir=tmp_path,
        files_dir=tmp_path / "files",
        artifacts_dir=tmp_path / "artifacts",
        cookie_secure=False,
    )

    def fake_provider(
        _provider_id: str, *, wants_artifact: bool, first_turn: bool
    ) -> MockProvider:
        del wants_artifact
        if first_turn:
            return MockProvider(
                text_chunks=("보고서를 구성하겠습니다.",),
                tool_call=MockToolCall(
                    name="create_report",
                    arguments=_arguments(report_format),
                    call_id=f"call_create_{report_format}",
                ),
            )
        return MockProvider(text_chunks=("보고서를 생성하고 형식을 검증했습니다.",))

    monkeypatch.setattr(local_run_executor, "_provider", fake_provider)
    with TestClient(create_app(settings)) as client:
        csrf = _login(client)
        project_id = client.get("/api/projects").json()[0]["id"]
        conversation = client.post(
            "/api/conversations",
            headers={"X-CSRF-Token": csrf},
            json={"projectId": project_id, "title": f"{report_format} 생성 테스트"},
        ).json()
        started = client.post(
            f"/api/conversations/{conversation['id']}/runs",
            headers={
                "X-CSRF-Token": csrf,
                "Idempotency-Key": f"create-{report_format}-report-0001",
            },
            json={
                "message": {
                    "text": f"점검 결과를 {report_format} 보고서 Artifact로 만들어 주세요.",
                    "attachmentIds": [],
                    "promptReferences": [],
                },
                "execution": {
                    "providerId": "mock",
                    "modelKey": "mock-agent",
                    "effortId": "medium",
                },
            },
        )
        assert started.status_code == 202, started.text
        snapshot = _wait_for_terminal(client, started.json()["run"]["runId"])
        assert snapshot["status"] == "completed"
        artifact = snapshot["artifacts"][0]
        expected_name, expected_kind, expected_mime = _EXPECTED_METADATA[report_format]
        assert artifact["displayName"] == expected_name
        assert artifact["kind"] == expected_kind
        assert artifact["mimeType"] == expected_mime
        assert artifact["validationStatus"] == (
            "passed" if report_format == "markdown" else "structural_passed"
        )

        version = client.get(f"/api/artifacts/{artifact['id']}/versions/1")
        assert version.status_code == 200


        version_payload = version.json()
        if report_format != "markdown":
            assert version_payload["validation"]["renderVerified"] is False
            assert version_payload["validation"]["renderer"] is None
            assert version_payload["validation"]["pages"] == []
            assert (
                "render_verification_pending"
                in version_payload["validation"]["warnings"]
            )
        if report_format in {"docx", "xlsx", "pptx", "pdf"}:
            assert version_payload["sourceText"] is None
            blocked_version = client.post(
                f"/api/artifacts/{artifact['id']}/versions",
                headers={
                    "X-CSRF-Token": csrf,
                    "If-Match": version_payload["etag"],
                    "Idempotency-Key": f"blocked-{report_format}-edit-0001",
                },
                json={
                    "baseVersion": 1,
                    "sourceText": "",
                    "changeSummary": "binary edit must be rejected",
                },
            )
            assert blocked_version.status_code == 409
            assert blocked_version.json()["code"] == "artifact_binary_edit_unsupported"
            blocked_draft = client.put(
                f"/api/artifacts/{artifact['id']}/draft",
                headers={"X-CSRF-Token": csrf},
                json={"baseVersion": 1, "content": ""},
            )
            assert blocked_draft.status_code == 409
            assert blocked_draft.json()["code"] == "artifact_binary_edit_unsupported"
        elif report_format == "markdown":
            assert version_payload["sourceText"] is not None
            draft = client.put(
                f"/api/artifacts/{artifact['id']}/draft",
                headers={"X-CSRF-Token": csrf},
                json={
                    "baseVersion": 1,
                    "content": version_payload["sourceText"] + "\n\n초안 수정",
                },
            )
            assert draft.status_code == 200, draft.text
            draft_files = list(
                (tmp_path / "artifacts" / "artifact-drafts").rglob("*.md")
            )
            assert len(draft_files) == 1
        downloaded = client.get(f"/api/artifacts/{artifact['id']}/download")
        assert downloaded.status_code == 200
        assert downloaded.headers["content-type"].split(";", 1)[0] == expected_mime
        assert downloaded.headers["content-disposition"].startswith(
            "attachment; filename*=UTF-8''%EA%B4%91%EC%96%91_"
        )
        assert downloaded.headers["content-disposition"].endswith(
            f".{expected_name.rsplit('.', 1)[-1]}"
        )
        _assert_reopened(report_format, downloaded.content)

        if report_format == "html":
            duplicate = client.post(
                f"/api/conversations/{conversation['id']}/runs",
                headers={
                    "X-CSRF-Token": csrf,
                    "Idempotency-Key": "create-html-report-duplicate-0001",
                },
                json={
                    "message": {
                        "text": "같은 주제의 HTML 보고서를 하나 더 만들어 주세요.",
                        "attachmentIds": [],
                        "promptReferences": [],
                    },
                    "execution": {
                        "providerId": "mock",
                        "modelKey": "mock-agent",
                        "effortId": "medium",
                    },
                },
            )
            assert duplicate.status_code == 202, duplicate.text
            duplicate_snapshot = _wait_for_terminal(
                client, duplicate.json()["run"]["runId"]
            )
            assert duplicate_snapshot["status"] == "completed"
            assert any(
                item["displayName"] == "광양_설비_점검_보고서_2.html"
                for item in duplicate_snapshot["artifacts"]
            )


def test_create_report_commit_failure_cleans_artifact_content(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = Settings(
        environment="test",
        database_url=f"sqlite:///{(tmp_path / 'report-cleanup.db').as_posix()}",
        data_dir=tmp_path,
        files_dir=tmp_path / "files",
        artifacts_dir=tmp_path / "artifacts",
        cookie_secure=False,
    )

    def fake_provider(
        _provider_id: str, *, wants_artifact: bool, first_turn: bool
    ) -> MockProvider:
        del wants_artifact
        if first_turn:
            return MockProvider(
                tool_call=MockToolCall(
                    name="create_report",
                    arguments=_arguments("markdown"),
                    call_id="call_report_cleanup",
                )
            )
        return MockProvider(text_chunks=("unexpected continuation",))

    real_commit = Session.commit
    artifact_commit_failed = False

    def fail_artifact_commit(session: Session) -> None:
        nonlocal artifact_commit_failed
        has_artifact_version = any(
            isinstance(item, ArtifactVersion) for item in session.identity_map.values()
        )
        if has_artifact_version and not artifact_commit_failed:
            artifact_commit_failed = True
            raise RuntimeError("forced create_report artifact commit failure")
        real_commit(session)

    monkeypatch.setattr(local_run_executor, "_provider", fake_provider)
    monkeypatch.setattr(Session, "commit", fail_artifact_commit)
    with TestClient(create_app(settings)) as client:
        csrf = _login(client)
        project_id = client.get("/api/projects").json()[0]["id"]
        conversation = client.post(
            "/api/conversations",
            headers={"X-CSRF-Token": csrf},
            json={"projectId": project_id, "title": "Report cleanup"},
        ).json()
        started = client.post(
            f"/api/conversations/{conversation['id']}/runs",
            headers={
                "X-CSRF-Token": csrf,
                "Idempotency-Key": "report-cleanup-0001",
            },
            json={"message": {"text": "Create a markdown report"}},
        )
        assert started.status_code == 202, started.text
        snapshot = _wait_for_terminal(client, started.json()["run"]["runId"])

    assert artifact_commit_failed is True
    assert snapshot["status"] == "failed"
    assert not [path for path in settings.artifacts_dir.rglob("*") if path.is_file()]


def test_selected_artifact_target_retries_short_report_and_exposes_separate_counts(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = Settings(
        environment="test",
        database_url=f"sqlite:///{(tmp_path / 'target-length.db').as_posix()}",
        data_dir=tmp_path,
        files_dir=tmp_path / "files",
        artifacts_dir=tmp_path / "artifacts",
        cookie_secure=False,
    )
    short_html = (
        "<!doctype html><html lang='ko'><head><title>짧은 보고서</title></head>"
        "<body><main><h1>짧은 보고서</h1><p>요약입니다.</p></main></body></html>"
    )
    first_extension = (
        "<section id='first-extension'><h2>첫 번째 추가 분석</h2>"
        + "<p>공급 구조와 계약 조건을 추가로 비교합니다.</p>\n" * 10
        + "</section>"
    )
    final_extension = (
        "<section id='final-extension'><h2>두 번째 추가 분석</h2>"
        + "<p>근거와 수치를 바탕으로 원인, 영향, 대응 방향을 구체적으로 분석합니다.</p>\n"
        * 300
        + "</section>"
    )
    provider_turn = 0
    requests = []

    class RecordingProvider(MockProvider):
        async def stream(self, request):
            requests.append(request)
            async for event in super().stream(request):
                yield event

    def fake_provider(
        _provider_id: str, *, wants_artifact: bool, first_turn: bool
    ) -> MockProvider:
        nonlocal provider_turn
        del first_turn
        assert wants_artifact is True
        provider_turn += 1
        if provider_turn == 1:
            arguments = _arguments("html")
            arguments["html_source"] = short_html
            return RecordingProvider(
                tool_call=MockToolCall(
                    name="create_report",
                    arguments=arguments,
                    call_id=f"call_target_length_{provider_turn}",
                )
            )
        if provider_turn <= 3:
            return RecordingProvider(
                tool_call=MockToolCall(
                    name="extend_report",
                    arguments={
                        "content": (
                            first_extension
                            if provider_turn == 2
                            else final_extension
                        )
                    },
                    call_id=f"call_target_length_{provider_turn}",
                )
            )
        return RecordingProvider(text_chunks=("확장한 HTML 보고서를 저장했습니다.",))

    monkeypatch.setattr(local_run_executor, "_provider", fake_provider)
    with TestClient(create_app(settings)) as client:
        csrf = _login(client)
        project_id = client.get("/api/projects").json()[0]["id"]
        conversation = client.post(
            "/api/conversations",
            headers={"X-CSRF-Token": csrf},
            json={"projectId": project_id, "title": "보고서 목표 분량"},
        ).json()
        started = client.post(
            f"/api/conversations/{conversation['id']}/runs",
            headers={
                "X-CSRF-Token": csrf,
                "Idempotency-Key": "target-length-report-0001",
            },
            json={
                "message": {
                    "text": "분석 보고서를 HTML로 만들어 주세요.",
                    "targetOutputTokens": 1_000,
                }
            },
        )
        assert started.status_code == 202, started.text
        snapshot = _wait_for_terminal(client, started.json()["run"]["runId"])
        artifact = client.get(
            f"/api/artifacts/{snapshot['artifacts'][0]['id']}"
        ).json()
        version_sources = [
            client.get(
                f"/api/artifacts/{snapshot['artifacts'][0]['id']}/versions/{version}"
            ).json()["sourceText"]
            for version in (1, 2, 3)
        ]

    assert snapshot["status"] == "completed"
    assert provider_turn == 4
    assert len(snapshot["artifacts"]) == 1
    assert snapshot["artifacts"][0]["currentVersion"] == 3
    assert artifact["versions"] == [3, 2, 1]
    assert len(snapshot["toolExecutions"]) == 3
    for attempt, execution in enumerate(snapshot["toolExecutions"][:2], start=1):
        result = execution["result"]
        assert result["status"] == "needs_expansion"
        assert result["artifact_id"] == snapshot["artifacts"][0]["id"]
        assert result["version"] == attempt
        assert result["documentTokens"] < result["minimumTokens"]
        assert result["expansionAttempt"] == attempt
        assert result["maxExpansionAttempts"] == 2
        assert execution["artifactId"] == snapshot["artifacts"][0]["id"]
    assert snapshot["toolExecutions"][2]["artifactId"] == snapshot["artifacts"][0]["id"]
    assert [
        execution["toolName"] for execution in snapshot["toolExecutions"]
    ] == ["create_report", "extend_report", "extend_report"]
    assert version_sources[0] == short_html
    assert short_html.split("</main>", 1)[0] in version_sources[1]
    assert first_extension in version_sources[1]
    assert short_html.split("</main>", 1)[0] in version_sources[2]
    assert first_extension in version_sources[2]
    assert final_extension in version_sources[2]
    assert version_sources[2].count("<!doctype html>") == 1
    artifact_usage = snapshot["artifactUsage"]
    assert artifact_usage["estimated"] is False
    assert artifact_usage["targetTokens"] == 1_000
    assert artifact_usage["tokens"] >= 800
    first_request = requests[0]
    system_text = "\n".join(
        str(message.content)
        for message in first_request.messages
        if message.role == "system"
    )
    assert "first-pass writing target" in system_text
    assert "acceptable first-call range is 80-105%" in system_text
    assert "about 800 to 1,050 tokens" in system_text
    assert "plan and draft near 90-100%—about 900 to 1,000 tokens" in system_text
    assert "Do not plan near the lower boundary" in system_text
    assert "start the `create_report` tool call before drafting the report body" in system_text
    assert "stream the complete report directly into its arguments" in system_text
    report_schema = next(
        schema
        for schema in first_request.tools
        if schema.get("function", {}).get("name") == "create_report"
    )
    assert (
        "selected document target is about 1,000 tokens"
        in report_schema["function"]["description"]
    )
    assert "Start this tool call as soon as research and analysis are ready" in (
        report_schema["function"]["description"]
    )
    assert "Prepare the complete first-pass report before calling this tool" not in (
        report_schema["function"]["description"]
    )
    html_description = report_schema["function"]["parameters"]["properties"][
        "html_source"
    ]["description"]
    assert "html_source itself must carry the full report content" in html_description
    assert "acceptable range is about 800 to 1,050 tokens" in html_description
    assert "prefer about 900 to 1,000 tokens" in html_description
    assert "at least 1,600 Unicode characters" in html_description
    assert (
        report_schema["function"]["parameters"]["properties"]["html_source"][
            "minLength"
        ]
        == 1_600
    )
    extension_schema = next(
        schema
        for schema in first_request.tools
        if schema.get("function", {}).get("name") == "extend_report"
    )
    assert "combines it with this fragment on the server" in (
        extension_schema["function"]["description"]
    )


def test_selected_artifact_target_fails_instead_of_saving_repeatedly_short_report(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = Settings(
        environment="test",
        database_url=f"sqlite:///{(tmp_path / 'target-length-failure.db').as_posix()}",
        data_dir=tmp_path,
        files_dir=tmp_path / "files",
        artifacts_dir=tmp_path / "artifacts",
        cookie_secure=False,
    )
    short_html = (
        "<!doctype html><html lang='ko'><head><title>짧은 보고서</title></head>"
        "<body><main><h1>짧은 보고서</h1><p>요약입니다.</p></main></body></html>"
    )
    extensions = (
        "<section id='failure-extension-1'><p>첫 번째 짧은 추가입니다.</p></section>",
        "<section id='failure-extension-2'><p>두 번째 짧은 추가입니다.</p></section>",
    )
    provider_turn = 0

    def fake_provider(
        _provider_id: str, *, wants_artifact: bool, first_turn: bool
    ) -> MockProvider:
        nonlocal provider_turn
        del first_turn
        assert wants_artifact is True
        provider_turn += 1
        if provider_turn == 1:
            arguments = _arguments("html")
            arguments["html_source"] = short_html
            tool_name = "create_report"
        else:
            arguments = {"content": extensions[provider_turn - 2]}
            tool_name = "extend_report"
        return MockProvider(
            tool_call=MockToolCall(
                name=tool_name,
                arguments=arguments,
                call_id=f"call_target_length_failure_{provider_turn}",
            )
        )

    monkeypatch.setattr(local_run_executor, "_provider", fake_provider)
    with TestClient(create_app(settings)) as client:
        csrf = _login(client)
        project_id = client.get("/api/projects").json()[0]["id"]
        conversation = client.post(
            "/api/conversations",
            headers={"X-CSRF-Token": csrf},
            json={"projectId": project_id, "title": "보고서 목표 분량 실패"},
        ).json()
        started = client.post(
            f"/api/conversations/{conversation['id']}/runs",
            headers={
                "X-CSRF-Token": csrf,
                "Idempotency-Key": "target-length-report-failure-0001",
            },
            json={
                "message": {
                    "text": "분석 보고서를 HTML로 만들어 주세요.",
                    "targetOutputTokens": 1_000,
                }
            },
        )
        assert started.status_code == 202, started.text
        snapshot = _wait_for_terminal(client, started.json()["run"]["runId"])
        final_source = client.get(
            f"/api/artifacts/{snapshot['artifacts'][0]['id']}/versions/3"
        ).json()["sourceText"]

    assert snapshot["status"] == "failed"
    assert snapshot["errorCode"] == "artifact_target_not_met"
    assert provider_turn == 3
    assert len(snapshot["artifacts"]) == 1
    assert snapshot["artifacts"][0]["currentVersion"] == 3
    assert len(snapshot["toolExecutions"]) == 3
    assert [
        execution["result"]["status"] for execution in snapshot["toolExecutions"][:2]
    ] == ["needs_expansion", "needs_expansion"]
    assert snapshot["toolExecutions"][2]["status"] == "failed"
    assert snapshot["toolExecutions"][2]["artifactId"] == snapshot["artifacts"][0]["id"]
    assert "최소 허용 분량" in snapshot["toolExecutions"][2]["error"]
    assert [
        execution["toolName"] for execution in snapshot["toolExecutions"]
    ] == ["create_report", "extend_report", "extend_report"]
    assert short_html.split("</main>", 1)[0] in final_source
    assert extensions[0] in final_source
    assert extensions[1] in final_source
    assert final_source.count("<!doctype html>") == 1


def test_short_report_rejects_full_document_retry_before_server_side_extension(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = Settings(
        environment="test",
        database_url=f"sqlite:///{(tmp_path / 'target-length-rewrite.db').as_posix()}",
        data_dir=tmp_path,
        files_dir=tmp_path / "files",
        artifacts_dir=tmp_path / "artifacts",
        cookie_secure=False,
    )
    short_html = (
        "<!doctype html><html lang='ko'><head><title>보존할 원본</title></head>"
        "<body><main><h1 id='original'>보존할 원본</h1></main></body></html>"
    )
    replacement_html = (
        "<!doctype html><html><head><title>덮어쓴 문서</title></head>"
        "<body><main><h1 id='replacement'>덮어쓴 문서</h1></main></body></html>"
    )
    extension = (
        "<section id='server-appended'><h2>누적 분석</h2>"
        + "<p>기존 근거를 유지하면서 원인과 영향을 추가로 분석합니다.</p>\n" * 300
        + "</section>"
    )
    provider_turn = 0

    def fake_provider(
        _provider_id: str, *, wants_artifact: bool, first_turn: bool
    ) -> MockProvider:
        nonlocal provider_turn
        del first_turn
        assert wants_artifact is True
        provider_turn += 1
        if provider_turn == 1:
            arguments = _arguments("html")
            arguments["html_source"] = short_html
            return MockProvider(
                tool_call=MockToolCall(
                    name="create_report",
                    arguments=arguments,
                    call_id="call_original_report",
                )
            )
        if provider_turn == 2:
            arguments = _arguments("html")
            arguments["html_source"] = replacement_html
            return MockProvider(
                tool_call=MockToolCall(
                    name="create_report",
                    arguments=arguments,
                    call_id="call_forbidden_rewrite",
                )
            )
        if provider_turn == 3:
            return MockProvider(
                tool_call=MockToolCall(
                    name="extend_report",
                    arguments={"content": extension},
                    call_id="call_server_extension",
                )
            )
        return MockProvider(text_chunks=("누적 확장한 보고서를 저장했습니다.",))

    monkeypatch.setattr(local_run_executor, "_provider", fake_provider)
    with TestClient(create_app(settings)) as client:
        csrf = _login(client)
        project_id = client.get("/api/projects").json()[0]["id"]
        conversation = client.post(
            "/api/conversations",
            headers={"X-CSRF-Token": csrf},
            json={"projectId": project_id, "title": "보고서 전체 재작성 차단"},
        ).json()
        started = client.post(
            f"/api/conversations/{conversation['id']}/runs",
            headers={
                "X-CSRF-Token": csrf,
                "Idempotency-Key": "target-length-rewrite-0001",
            },
            json={
                "message": {
                    "text": "분석 보고서를 HTML로 만들어 주세요.",
                    "targetOutputTokens": 1_000,
                }
            },
        )
        assert started.status_code == 202, started.text
        snapshot = _wait_for_terminal(client, started.json()["run"]["runId"])
        artifact_id = snapshot["artifacts"][0]["id"]
        artifact = client.get(f"/api/artifacts/{artifact_id}").json()
        final_source = client.get(
            f"/api/artifacts/{artifact_id}/versions/2"
        ).json()["sourceText"]

    assert snapshot["status"] == "completed"
    assert provider_turn == 4
    assert artifact["versions"] == [2, 1]
    assert [
        execution["toolName"] for execution in snapshot["toolExecutions"]
    ] == ["create_report", "create_report", "extend_report"]
    rejected = snapshot["toolExecutions"][1]
    assert rejected["status"] == "failed"
    assert "`extend_report`" in rejected["error"]
    assert "id='original'" in final_source
    assert "id='server-appended'" in final_source
    assert "id='replacement'" not in final_source


def _assert_reopened(report_format: str, content: bytes) -> None:
    if report_format in {"html", "markdown"}:
        source = content.decode("utf-8", errors="strict")
        assert "광양 설비 점검 보고서" in source
        assert "이상 징후 1건" in source
        return
    if report_format == "docx":
        document = Document(BytesIO(content))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        text += "\n".join(
            cell.text
            for table in document.tables
            for row in table.rows
            for cell in row.cells
        )
        assert "광양 설비 점검 보고서" in text
        assert "이상 징후 1건" in text
        return
    if report_format == "xlsx":
        workbook = load_workbook(BytesIO(content), data_only=False)
        values = "\n".join(
            str(cell.value)
            for sheet in workbook.worksheets
            for row in sheet.iter_rows()
            for cell in row
            if cell.value is not None
        )
        assert "광양 설비 점검 보고서" in values
        assert "이상 징후 1건" in values
        workbook.close()
        return
    if report_format == "pptx":
        presentation = Presentation(BytesIO(content))
        text = "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert "광양 설비 점검 보고서" in text
        assert "이상 징후 1건" in text
        return
    if report_format == "pdf":
        reader = PdfReader(BytesIO(content))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        assert "광양 설비 점검 보고서" in text
        assert "이상 징후 1건" in text
        assert float(reader.pages[0].mediabox.width) == pytest.approx(595.28, abs=0.5)
        assert float(reader.pages[0].mediabox.height) == pytest.approx(841.89, abs=0.5)
        return
    raise AssertionError(f"Unhandled report format: {report_format}")


def _login(client: TestClient) -> str:
    response = client.post(
        "/api/auth/login",
        json={
            "loginName": "admin",
            "loginDomain": "posco.com",
            "password": "1",
        },
    )
    assert response.status_code == 200
    return response.json()["csrfToken"]


def _wait_for_terminal(client: TestClient, run_id: str) -> dict[str, object]:
    deadline = time.monotonic() + 10
    while time.monotonic() < deadline:
        response = client.get(f"/api/runs/{run_id}/snapshot")
        assert response.status_code == 200
        payload = response.json()
        if payload["status"] in {
            "completed",
            "failed",
            "cancelled",
            "limit_reached",
            "interrupted",
        }:
            return payload
        time.sleep(0.03)
    raise AssertionError("Run did not reach a terminal state")
