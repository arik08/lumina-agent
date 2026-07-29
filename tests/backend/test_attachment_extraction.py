from __future__ import annotations

import asyncio
import threading
import time
from io import BytesIO
from pathlib import Path
from unittest.mock import patch

import pytest
from docx import Document
from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter
from sqlalchemy.orm import Session

from lumina.api.routes.attachments import _extract_attachment, _sniff_mime
from lumina.attachments import extract_attachment_text
from lumina.config import Settings
from lumina.main import create_app


def test_attachment_extraction_runs_off_the_event_loop_thread() -> None:
    caller_thread = threading.get_ident()
    worker_threads: list[int] = []

    def recording_extractor(**kwargs):
        worker_threads.append(threading.get_ident())
        return extract_attachment_text(**kwargs)

    async def extract():
        with patch(
            "lumina.api.routes.attachments.extract_attachment_text",
            side_effect=recording_extractor,
        ):
            return await _extract_attachment(
                filename="inspection.txt",
                mime_type="text/plain",
                content=b"inspection",
            )

    result = asyncio.run(extract())

    assert result.status == "completed"
    assert worker_threads and worker_threads[0] != caller_thread


def test_attachment_extraction_worker_concurrency_is_bounded() -> None:
    active = 0
    peak = 0
    lock = threading.Lock()

    def recording_extractor(**kwargs):
        nonlocal active, peak
        with lock:
            active += 1
            peak = max(peak, active)
        try:
            time.sleep(0.03)
            return extract_attachment_text(**kwargs)
        finally:
            with lock:
                active -= 1

    async def extract_all():
        with patch(
            "lumina.api.routes.attachments.extract_attachment_text",
            side_effect=recording_extractor,
        ):
            return await asyncio.gather(
                *(
                    _extract_attachment(
                        filename=f"inspection-{index}.txt",
                        mime_type="text/plain",
                        content=b"inspection",
                    )
                    for index in range(6)
                )
            )

    results = asyncio.run(extract_all())

    assert all(result.status == "completed" for result in results)
    assert peak == 2


def test_text_extraction_tracks_lines() -> None:
    result = extract_attachment_text(
        filename="inspection.md",
        mime_type="text/markdown",
        content="첫 줄\n둘째 줄".encode(),
    )
    assert result.status == "completed"
    assert result.text == "첫 줄\n둘째 줄"
    assert result.locator_map == {"kind": "line", "count": 2}


def test_pdf_extraction_allows_public_permission_encryption() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    writer.encrypt(user_password="", owner_password="owner-secret")
    encrypted_buffer = BytesIO()
    writer.write(encrypted_buffer)

    extracted = extract_attachment_text(
        filename="public-report.pdf",
        mime_type="application/pdf",
        content=encrypted_buffer.getvalue(),
    )

    assert extracted.status == "completed"
    assert "[Page 1]" in extracted.text


def test_pdf_extraction_rejects_required_reader_password() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    writer.encrypt(user_password="reader-secret", owner_password="owner-secret")
    encrypted_buffer = BytesIO()
    writer.write(encrypted_buffer)

    extracted = extract_attachment_text(
        filename="private-report.pdf",
        mime_type="application/pdf",
        content=encrypted_buffer.getvalue(),
    )

    assert extracted.status == "failed"
    assert extracted.metadata["errorType"] == "EncryptedPdf"


def test_office_formats_extract_without_writing_temporary_files() -> None:
    document = Document()
    document.add_heading("설비 점검", level=1)
    document.add_paragraph("베어링 온도를 확인했습니다.")
    docx_buffer = BytesIO()
    document.save(docx_buffer)
    docx = extract_attachment_text(
        filename="inspection.docx",
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        content=docx_buffer.getvalue(),
    )
    assert docx.status == "completed"
    assert "베어링 온도" in docx.text

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "점검"
    sheet.append(["설비", "상태"])
    sheet.append(["압연기", "정상"])
    xlsx_buffer = BytesIO()
    workbook.save(xlsx_buffer)
    workbook.close()
    xlsx = extract_attachment_text(
        filename="inspection.xlsx",
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        content=xlsx_buffer.getvalue(),
    )
    assert xlsx.status == "completed"
    assert "[Sheet: 점검]" in xlsx.text
    assert "압연기\t정상" in xlsx.text

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "2분기 결과"
    slide.placeholders[1].text = "예방 정비 완료"
    pptx_buffer = BytesIO()
    presentation.save(pptx_buffer)
    pptx = extract_attachment_text(
        filename="inspection.pptx",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        content=pptx_buffer.getvalue(),
    )
    assert pptx.status == "completed"
    assert "예방 정비 완료" in pptx.text

    assert _sniff_mime(docx_buffer.getvalue(), ".docx").endswith(
        "wordprocessingml.document"
    )
    assert _sniff_mime(xlsx_buffer.getvalue(), ".xlsx").endswith("spreadsheetml.sheet")
    assert _sniff_mime(pptx_buffer.getvalue(), ".pptx").endswith(
        "presentationml.presentation"
    )


def test_zip_header_alone_is_not_trusted_as_an_office_document() -> None:
    fake_zip = b"PK\x03\x04not-an-openxml-container"

    assert _sniff_mime(fake_zip, ".docx") == "application/octet-stream"
    assert _sniff_mime(fake_zip, ".xlsx") == "application/octet-stream"
    assert _sniff_mime(fake_zip, ".pptx") == "application/octet-stream"


def test_attachment_api_rejects_fake_office_and_persists_valid_extraction(
    tmp_path: Path,
) -> None:
    settings = Settings(
        environment="test",
        database_url=f"sqlite:///{(tmp_path / 'attachments.db').as_posix()}",
        data_dir=tmp_path,
        files_dir=tmp_path / "files",
        artifacts_dir=tmp_path / "artifacts",
        cookie_secure=False,
    )
    with TestClient(create_app(settings)) as client:
        login = client.post(
            "/api/auth/login",
            json={
                "loginName": "admin",
                "loginDomain": "posco.com",
                "password": "1111",
            },
        )
        csrf = login.json()["csrfToken"]
        project_id = client.get("/api/projects").json()[0]["id"]
        conversation = client.post(
            "/api/conversations",
            headers={"X-CSRF-Token": csrf},
            json={"projectId": project_id, "title": "첨부 검증"},
        ).json()
        endpoint = f"/api/conversations/{conversation['id']}/attachments"

        rejected = client.post(
            endpoint,
            headers={"X-CSRF-Token": csrf},
            files={
                "file": (
                    "fake.docx",
                    b"PK\x03\x04not-openxml",
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
        )
        assert rejected.status_code == 415
        assert rejected.json()["code"] == "mime_mismatch"

        empty = client.post(
            endpoint,
            headers={"X-CSRF-Token": csrf},
            data={"pasted_text": "  \n  ", "source": "paste"},
        )
        assert empty.status_code == 422
        assert empty.json()["code"] == "attachment_empty"

        document = Document()
        document.add_paragraph("실제 문서 본문")
        buffer = BytesIO()
        document.save(buffer)
        uploaded = client.post(
            endpoint,
            headers={"X-CSRF-Token": csrf},
            files={
                "file": (
                    "real.docx",
                    buffer.getvalue(),
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
        )
        assert uploaded.status_code == 201, uploaded.text
        assert uploaded.json()["extractionStatus"] == "completed"
        assert uploaded.json()["metadata"]["extractedSize"] > 0

        stored_attachment = next(
            path
            for path in (settings.files_dir / "attachments").rglob("*")
            if path.is_file()
        )
        stored_attachment.write_bytes(b"tampered attachment")
        unavailable = client.get(
            f"/api/attachments/{uploaded.json()['id']}/content"
        )
        assert unavailable.status_code == 503
        assert unavailable.json()["code"] == "attachment_content_missing"


def test_attachment_commit_failure_cleans_all_managed_files(tmp_path: Path) -> None:
    settings = Settings(
        environment="test",
        database_url=f"sqlite:///{(tmp_path / 'attachment-cleanup.db').as_posix()}",
        data_dir=tmp_path,
        files_dir=tmp_path / "files",
        artifacts_dir=tmp_path / "artifacts",
        cookie_secure=False,
    )
    with TestClient(create_app(settings)) as client:
        login = client.post(
            "/api/auth/login",
            json={
                "loginName": "admin",
                "loginDomain": "posco.com",
                "password": "1111",
            },
        )
        csrf = login.json()["csrfToken"]
        project_id = client.get("/api/projects").json()[0]["id"]
        conversation = client.post(
            "/api/conversations",
            headers={"X-CSRF-Token": csrf},
            json={"projectId": project_id, "title": "첨부 rollback 검증"},
        ).json()

        with patch.object(
            Session, "commit", side_effect=RuntimeError("forced attachment commit failure")
        ):
            with pytest.raises(RuntimeError, match="forced attachment commit failure"):
                client.post(
                    f"/api/conversations/{conversation['id']}/attachments",
                    headers={"X-CSRF-Token": csrf},
                    data={"pasted_text": "commit 실패 시 정리할 본문", "source": "paste"},
                )

    assert not [path for path in settings.files_dir.rglob("*") if path.is_file()]
